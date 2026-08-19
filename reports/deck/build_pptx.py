"""Build the DCS 404 deck as a PowerPoint file that Canva can import and edit.

Everything here is a native PowerPoint shape or text box -- no flattened images,
no outlined text -- so every headline, label and bar can be selected and changed
in Canva. Fonts are named, not embedded, which is the whole point: Canva matches
by name, so both faces are ones its library carries.
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---- the two faces, chosen because Canva's library carries both -------------
SERIF = "Noto Serif"
MONO  = "IBM Plex Mono"

# ---- palette, carried over from the HTML deck ------------------------------
PAPER   = RGBColor(0xFA, 0xFA, 0xF7)
INK     = RGBColor(0x1B, 0x1B, 0x18)
MUTED   = RGBColor(0x6B, 0x6A, 0x62)
HAIR    = RGBColor(0xDC, 0xDA, 0xD2)
HAIR2   = RGBColor(0xEA, 0xE8, 0xE1)
BLUE    = RGBColor(0x2E, 0x4F, 0xA8)
BLUE_LT = RGBColor(0x93, 0xA9, 0xDC)
CRIMSON = RGBColor(0xB3, 0x38, 0x2C)
TEAL    = RGBColor(0x2F, 0x6F, 0x62)

# The HTML stage is 1280x720 css px; a 16:9 slide is 13.333x7.5in = 1280x720 @96dpi,
# so one css pixel maps to exactly one slide pixel and the layout transfers 1:1.
def X(px): return Emu(int(px * 914400 / 96))
def P(px): return Pt(px * 0.75)

PAD_L, PAD_T, PAD_R = 96, 78, 96
CONTENT_W = 1280 - PAD_L - PAD_R


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    return s


def text(slide, x, y, w, runs, size=15, font=SERIF, color=INK, bold=False,
         italic=False, align=PP_ALIGN.LEFT, spc=None, line=1.35, h=60, caps=False):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    box = slide.shapes.add_textbox(X(x), X(y), X(w), X(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for t, ov in runs:
        r = p.add_run()
        r.text = t.upper() if caps else t
        f = r.font
        f.size = P(ov.get("size", size))
        f.name = ov.get("font", font)
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.color.rgb = ov.get("color", color)
        s = ov.get("spc", spc)
        if s:
            r._r.get_or_add_rPr().set("spc", str(int(s * 100)))
    return box


def rect(slide, x, y, w, h, fill=None, line_col=None, line_w=1, dash=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, X(x), X(y), X(w), X(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_col is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_col; sh.line.width = Pt(line_w)
        if dash:
            sh.line.dash_style = dash
    sh.text_frame.text = ""
    return sh


def oval(slide, cx, cy, r, fill=None, line_col=None, line_w=1.2):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, X(cx - r), X(cy - r), X(2 * r), X(2 * r))
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_col is None: sh.line.fill.background()
    else: sh.line.color.rgb = line_col; sh.line.width = Pt(line_w)
    return sh


def seg(slide, x1, y1, x2, y2, color=HAIR, w=1):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, X(x1), X(y1), X(x2), X(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    c.shadow.inherit = False
    return c


def hrule(slide, x, y, w, color=HAIR, weight=1):
    seg(slide, x, y, x + w, y, color, weight)


def head(slide, sec, label, title, dek=None, title_w=None):
    """The eyebrow + headline + optional single dek line every content slide opens with."""
    text(slide, PAD_L, PAD_T, CONTENT_W,
         [(f"§{sec}", {"color": CRIMSON}), ("   " + label.upper(), {"color": MUTED})],
         size=11, font=MONO, spc=1.6, h=20)
    text(slide, PAD_L, PAD_T + 22, title_w or 820, title, size=46, bold=True,
         line=1.1, h=130)
    y = PAD_T + 22 + (130 if "\n" in title else 62)
    if dek:
        text(slide, PAD_L, PAD_T + 96, 700, dek, size=20, italic=True,
             color=MUTED, line=1.45, h=70)
    return y


def head(slide, sec, label, title, tl=1, dek=None, dl=1, dw=700, tw=860):
    """Eyebrow, headline, optional dek. Returns the y where slide content starts."""
    text(slide, PAD_L, PAD_T, CONTENT_W,
         [(f"§{sec}", {"color": CRIMSON}), ("   " + label.upper(), {"color": MUTED})],
         size=11, font=MONO, spc=1.6, h=20)
    text(slide, PAD_L, PAD_T + 22, tw, title, size=46, bold=True, line=1.1, h=60 * tl + 20)
    y = PAD_T + 22 + (56 if tl == 1 else 108)
    if dek:
        text(slide, PAD_L, y + 14, dw, dek, size=20, italic=True, color=MUTED,
             line=1.45, h=34 * dl + 10)
        y += 14 + 30 * dl
    return y + 34


def stat_band(slide, y, cells, h=118):
    """The rule-bounded row of figures used on the problem and dataset slides."""
    hrule(slide, PAD_L, y, CONTENT_W, INK, 1.2)
    w = CONTENT_W / len(cells)
    for i, (val, key, col) in enumerate(cells):
        x = PAD_L + i * w
        if i:
            seg(slide, x, y + 8, x, y + h - 8, HAIR, 1)
        text(slide, x + (0 if i == 0 else 24), y + 18, w - 30, val,
             size=38, bold=True, color=col, h=52)
        text(slide, x + (0 if i == 0 else 24), y + 78, w - 30, key,
             size=10.5, font=MONO, color=MUTED, spc=1.3, line=1.6, h=32, caps=True)
    hrule(slide, PAD_L, y + h, CONTENT_W, HAIR, 1)


def kstats(slide, x, y, w, rows, vh=104):
    """The stacked figures down the right of the results and GNN slides."""
    hrule(slide, x, y, w, INK, 1.2)
    for i, (val, key, col) in enumerate(rows):
        yy = y + i * vh
        text(slide, x, yy + 16, w, val, size=40, bold=True, color=col, h=46)
        text(slide, x, yy + 72, w, key, size=10.5, font=MONO, color=MUTED,
             spc=1.2, line=1.6, h=28, caps=True)
        hrule(slide, x, yy + vh, w, HAIR, 1)


def barchart(slide, x, y, w, groups, tick_lab=True, marker=None, LBL=170):
    """Horizontal bars drawn as real rectangles so every bar stays editable."""
    x0, x1 = x + LBL, x + w
    bh, gin, ggap, top = 28, 6, 40, y
    def SX(v): return x0 + (x1 - x0) * v
    nb = sum(len(g[1]) for g in groups)
    H = nb * bh + sum((len(g[1]) - 1) * gin for g in groups) + (len(groups) - 1) * ggap
    bot = top + H
    for t in (0, .25, .5, .75, 1):
        seg(slide, SX(t), top - 12, SX(t), bot + 8, HAIR2, 1)
        if tick_lab:
            text(slide, SX(t) - 30, bot + 14, 60, f"{t:.2f}", size=11, font=MONO,
                 color=MUTED, align=PP_ALIGN.CENTER, h=18)
    yy = top
    for gname, rows in groups:
        gh = len(rows) * bh + (len(rows) - 1) * gin
        text(slide, x, yy + gh / 2 - 9, LBL - 16, gname, size=12, font=MONO,
             align=PP_ALIGN.RIGHT, h=20)
        for lab, v, col in rows:
            bw = SX(v) - x0
            rect(slide, x0, yy, bw, bh, fill=col)
            text(slide, x0 + bw + 10, yy + bh / 2 - 9, 90, f"{v:.3f}",
                 size=13, font=MONO, bold=True, h=20)
            yy += bh + gin
        yy += ggap - gin
    seg(slide, x0, top - 12, x0, bot + 8, INK, 1.2)
    if marker is not None:
        mv, mlab = marker
        mk = seg(slide, SX(mv), top - 18, SX(mv), bot + 8, CRIMSON, 1.4)
        mk.line.dash_style = DASH.DASH
        text(slide, SX(mv) - 160, top - 38, 155, mlab, size=10.5, font=MONO,
             color=CRIMSON, align=PP_ALIGN.RIGHT, h=16)
    return bot


def bell(slide, cx, base, wid, hgt, fill, line_col):
    """A distribution hump, approximated with short segments so it stays a shape."""
    import math
    pts = []
    n = 26
    for i in range(n + 1):
        t = -1 + 2 * i / n
        px = cx + t * wid / 2
        py = base - hgt * math.exp(-4.0 * t * t)
        pts.append((X(px), X(py)))
    b = slide.shapes.build_freeform(X(cx - wid / 2), X(base))
    b.add_line_segments(pts + [(X(cx + wid / 2), X(base))], close=True)
    sh = b.convert_to_shape()
    sh.shadow.inherit = False
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line_col; sh.line.width = Pt(1.4)
    return sh


from pptx.enum.dml import MSO_LINE_DASH_STYLE as DASH

BELL_C = RGBColor(0xEE, 0xDB, 0xD6)   # crimson at 16% over paper
BELL_B = RGBColor(0xDD, 0xE2, 0xEC)   # blue at 14% over paper
TRAIN  = RGBColor(0xC6, 0xD1, 0xE9)
TEST   = RGBColor(0xF0, 0xDF, 0xDA)

prs = Presentation()
prs.slide_width, prs.slide_height = X(1280), X(720)

# ============================ 1 - TITLE =====================================
s = new_slide(prs)
text(s, PAD_L, PAD_T, 500, "DCS 404 · FINAL PROJECT", size=12, font=MONO,
     color=CRIMSON, spc=1.8, h=20)
text(s, PAD_L, 96, 690,
     "Detecting Illicit Ethereum Transactions with Graph‑Derived Features",
     size=48, bold=True, line=1.1, h=300)
rect(s, PAD_L, 404, 2, 84, fill=BLUE)
text(s, PAD_L + 26, 406, 470,
     "Can an address’s network of connections reveal fraud that its own features cannot?",
     size=22, italic=True, color=BLUE, line=1.45, h=80)
text(s, PAD_L, 518, 520, [("Sanskriti Acharya", {}), (" & ", {"color": MUTED, "italic": True}),
                          ("Sweta Sharma", {})], size=20, h=30)
text(s, PAD_L, 552, 520, "eth-illicit-detection.streamlit.app", size=12, font=MONO,
     color=MUTED, spc=0.6, h=20)

GX, GY, GS = 762, 168, 1.02
def gx(v): return GX + v * GS
def gy(v): return GY + v * GS
for a, b in [((180,160),(82,84)), ((180,160),(262,70)), ((180,160),(300,176)),
             ((180,160),(236,268)), ((180,160),(96,238)), ((82,84),(30,140)),
             ((82,84),(140,30)), ((262,70),(206,24)), ((262,70),(330,102)),
             ((300,176),(330,102)), ((96,238),(42,288)), ((96,238),(30,140)),
             ((236,268),(300,300))]:
    seg(s, gx(a[0]), gy(a[1]), gx(b[0]), gy(b[1]), HAIR, 1.2)
for c in [(82,84),(262,70),(300,176),(236,268),(96,238)]:
    oval(s, gx(c[0]), gy(c[1]), 9 * GS, fill=PAPER, line_col=MUTED, line_w=1.4)
for c in [(30,140),(140,30),(206,24),(330,102),(42,288),(300,300)]:
    oval(s, gx(c[0]), gy(c[1]), 6 * GS, fill=PAPER, line_col=MUTED, line_w=1.4)
ring = oval(s, gx(180), gy(160), 20 * GS, fill=None, line_col=CRIMSON, line_w=1)
ring.line.dash_style = DASH.DASH
oval(s, gx(180), gy(160), 13 * GS, fill=CRIMSON)
text(s, gx(180) - 90, gy(207), 180, "0x3f4e…a71c", size=11, font=MONO,
     color=CRIMSON, align=PP_ALIGN.CENTER, spc=1, h=18)

# ============================ 2 - PROBLEM ===================================
s = new_slide(prs)
y = head(s, 1, "Problem", "Is this address a scam?",
         dek="Every Ethereum transfer is public and permanent — and there are over a "
             "million a day. No investigator can read them all.", dl=2)
stat_band(s, y, [("Track A", "supervised · binary", INK),
                 ("1 address", "in · phishing or scam?", INK),
                 ("a probability", "out · never a verdict", BLUE)])
text(s, PAD_L, y + 164, 780,
     "Built for an analyst triaging addresses, so the output is a review priority. "
     "The deliverable is a short, ranked list.", size=20, italic=True, color=MUTED,
     line=1.45, h=70)

# ============================ 3 - METRIC ====================================
s = new_slide(prs)
y = head(s, 1, "Metric", "Accuracy is unusable here")
hrule(s, PAD_L, y, CONTENT_W, INK, 1.2)
for i, (v, k) in enumerate([("0.845", "accuracy of a model that always answers “licit”"),
                            ("F₁ = 0.000", "the same model, on the metric that matters")]):
    x = PAD_L + i * CONTENT_W / 2
    if i: seg(s, x, y + 8, x, y + 132, HAIR, 1)
    text(s, x + (0 if i == 0 else 30), y + 22, CONTENT_W / 2 - 40, v, size=56,
         bold=True, color=CRIMSON, h=76)
    text(s, x + (0 if i == 0 else 30), y + 106, CONTENT_W / 2 - 40, k, size=11,
         font=MONO, color=MUTED, spc=1.1, line=1.7, h=34, caps=True)
hrule(s, PAD_L, y + 142, CONTENT_W, HAIR, 1)
text(s, PAD_L, y + 180, 780, "F₁ = 2PR / (P + R)  ·  on the illicit class, reported with PR‑AUC",
     size=24, italic=True, color=MUTED, h=36)
text(s, PAD_L, y + 232, 800,
     "A missed scam lets fraud continue; a false alarm costs an analyst minutes. "
     "Recall matters more — but F₁ only rises when both are respectable.",
     size=20, italic=True, color=MUTED, line=1.45, h=70)

# ============================ 4 - DATASET ===================================
s = new_slide(prs)
y = head(s, 2, "Dataset", "Only the labels are borrowed")
stat_band(s, y, [("2,852", "labelled addresses", INK), ("652", "illicit · 22.9%", CRIMSON),
                 ("2017–18", "both classes, same era", INK), ("100%", "features computed by us", BLUE)])
for i, (lead, body) in enumerate([
        ("Labels", " — the MyEtherWallet darklist of reported scam addresses, deduplicated 715 → 652 wallets."),
        ("Everything else", " — raw transactions fetched address-by-address from Blockscout. No pre-cleaned download exists.")]):
    text(s, PAD_L + i * (CONTENT_W / 2 + 18), y + 164, CONTENT_W / 2 - 36,
         [(lead, {"bold": True, "color": INK, "italic": False}), (body, {})],
         size=20, italic=True, color=MUTED, line=1.45, h=110)

# ============================ 5 - THE TRAP ==================================
s = new_slide(prs)
y = head(s, 2, "Sampling", "The shortcut we had to close",
         dek="97% of the darklist dates to 2017–18. Sample licit addresses from recent "
             "blocks and block_number alone separates the classes perfectly.", dl=2, dw=820)
for i, (lab, ok) in enumerate([("NAIVE: LICIT FROM RECENT BLOCKS", False),
                               ("ERA-MATCHED (OURS)", True)]):
    bx = PAD_L + i * 560
    text(s, bx, y, 520, lab, size=10.5, font=MONO, color=MUTED, spc=1.5, h=18)
    base = y + 132
    if ok:
        bell(s, bx + 150, base, 250, 92, BELL_C, CRIMSON)
        bell(s, bx + 210, base, 250, 88, BELL_B, BLUE)
    else:
        bell(s, bx + 120, base, 250, 92, BELL_C, CRIMSON)
        bell(s, bx + 380, base, 200, 84, BELL_B, BLUE)
        text(s, bx + 80, base + 8, 120, "illicit", size=10, font=MONO, color=CRIMSON,
             align=PP_ALIGN.CENTER, h=16)
        text(s, bx + 330, base + 8, 120, "licit", size=10, font=MONO, color=BLUE,
             align=PP_ALIGN.CENTER, h=16)
    hrule(s, bx, base, 500, INK, 1)
    text(s, bx, base + 34, 520,
         "separable by block number alone ✗" if not ok else "distributions overlap — shortcut closed ✓",
         size=14, italic=True, color=CRIMSON if not ok else INK, h=22)
text(s, PAD_L, y + 210, CONTENT_W, "block number →", size=10, font=MONO, color=MUTED,
     align=PP_ALIGN.CENTER, spc=1, h=16)
text(s, PAD_L, y + 244, CONTENT_W,
     [("Fig. 1", {"bold": True, "color": INK}),
      (" — same logic twice more: contracts excluded from both classes, absolute block numbers excluded from features.", {})],
     size=11.5, font=MONO, color=MUTED, align=PP_ALIGN.CENTER, line=1.75, h=40)

# ============================ 6 - FEATURES ==================================
s = new_slide(prs)
y = head(s, 3, "EDA & Preprocessing", "29 behavioural features,\nplus 7 from the graph", tl=2, tw=640)
LX = PAD_L
text(s, LX, y, 300, "value_concentration ≈ 1", size=10.5, font=MONO, color=MUTED, spc=1, h=18)
for vy in (28, 56, 84, 110):
    seg(s, LX + 26, y + vy, LX + 140, y + 68, MUTED, 1.1)
for vy in (26, 54, 84, 112):
    oval(s, LX + 22, y + vy, 6, fill=PAPER, line_col=MUTED, line_w=1.3)
oval(s, LX + 152, y + 68, 11, fill=CRIMSON)
rect(s, LX + 164, y + 65, 86, 5, fill=CRIMSON)
oval(s, LX + 268, y + 68, 8, fill=PAPER, line_col=CRIMSON, line_w=1.6)
text(s, LX, y + 140, 400, "many victims in — one sweep out", size=14.5, italic=True, h=24)
RX = PAD_L + 330
text(s, RX, y, 300, "burstiness ≫ 0", size=10.5, font=MONO, color=MUTED, spc=1, h=18)
hrule(s, RX, y + 46, 226, HAIR, 1)
for i in range(8):
    seg(s, RX + 12 + i * 28, y + 38, RX + 12 + i * 28, y + 54, BLUE, 1.6)
text(s, RX, y + 62, 320, "ordinary wallet — steady", size=14, italic=True, color=BLUE, h=22)
hrule(s, RX, y + 108, 226, HAIR, 1)
for i, off in enumerate([8, 14, 20, 27, 33, 41]):
    seg(s, RX + off, y + 100, RX + off, y + 116, CRIMSON, 1.6)
text(s, RX, y + 124, 320, "scam wallet — burst, then silence", size=14.5, italic=True, color=CRIMSON, h=22)
text(s, PAD_L, y + 184, 560,
     [("Fig. 2", {"bold": True, "color": INK}),
      (" — two features that encode a fraud pattern. Graph features: degree, in/out degree, "
       "PageRank, clustering, community size, neighbour risk.", {})],
     size=11.5, font=MONO, color=MUTED, line=1.75, h=60)
TX = PAD_L + 620
rows = [("715 → 652 labels", "one wallet, several campaigns — duplicates would span the split"),
        ("empty addresses kept", "deleting them shifts class balance and hides a real case"),
        ("gasUsed dropped", "the API returns 0 on many real transfers"),
        ("scaler in pipeline", "the app cannot preprocess differently than training did")]
hrule(s, TX, y, 468, INK, 1)
ry = y
for k, v in rows:
    text(s, TX, ry + 16, 180, k, size=13, font=MONO, h=40)
    text(s, TX + 196, ry + 14, 272, v, size=17, italic=True, color=MUTED, line=1.4, h=58)
    ry += 76
    hrule(s, TX, ry, 468, HAIR2, 1)

# ============================ 7 - THE SPLIT =================================
s = new_slide(prs)
y = head(s, 4, "Modeling", "The split is temporal, not random",
         dek="Scam operations run clusters of sibling wallets. A random split puts siblings "
             "on both sides, and the model recognises the cluster instead of generalising.",
         dl=2, dw=840)
BW = CONTENT_W
rect(s, PAD_L, y + 20, BW * 0.7, 54, fill=TRAIN, line_col=BLUE, line_w=1.2)
rect(s, PAD_L + BW * 0.7, y + 20, BW * 0.3, 54, fill=TEST, line_col=CRIMSON, line_w=1.2)
text(s, PAD_L, y + 38, BW * 0.7, "TRAIN · earliest 70%", size=12, font=MONO,
     align=PP_ALIGN.CENTER, h=20)
text(s, PAD_L + BW * 0.7, y + 38, BW * 0.3, "TEST · latest 30%", size=12, font=MONO,
     color=CRIMSON, align=PP_ALIGN.CENTER, h=20)
text(s, PAD_L, y + 92, 400, "5-fold CV & tuning here", size=10.5, font=MONO, color=MUTED, h=18)
text(s, PAD_L + BW - 400, y + 92, 400, "touched once, at scoring", size=10.5, font=MONO,
     color=CRIMSON, align=PP_ALIGN.RIGHT, h=18)
text(s, PAD_L, y + 140, CONTENT_W,
     [("Fig. 3", {"bold": True, "color": INK}),
      (" — addresses ordered by first on-chain appearance. Learn from what is known; score what came later.", {})],
     size=11.5, font=MONO, color=MUTED, align=PP_ALIGN.CENTER, line=1.75, h=40)

# ============================ 8 - LEAKAGE ===================================
s = new_slide(prs)
y = head(s, 4, "Leakage", "The feature that could have leaked")
rect(s, PAD_L, y + 6, 2, 112, fill=CRIMSON)
text(s, PAD_L + 28, y + 6, 620,
     [("neighbour_risk_ratio", {"font": MONO, "size": 17}),
      (" asks “what fraction of my labelled neighbours are illicit?” Computed over ", {}),
      ("all", {"italic": True}), (" labels, it leaks the answer outright.", {})],
     size=21, line=1.5, h=110)
rect(s, PAD_L, y + 138, 2, 66, fill=BLUE)
text(s, PAD_L + 28, y + 140, 620,
     [("So ", {}), ("graph.py", {"font": MONO, "size": 17}),
      (" takes visible labels as an argument, and training passes ", {}),
      ("the training split’s labels only", {"bold": True}), (".", {})],
     size=21, line=1.5, h=80)
text(s, PAD_L, y + 242, 900,
     [("PINNED BY UNIT TESTS — ", {}),
      ("INCLUDING THAT AN ADDRESS CAN NEVER READ ITS OWN LABEL BACK THROUGH A SELF-LOOP.", {"color": INK})],
     size=12, font=MONO, color=MUTED, spc=0.5, line=1.8, h=44)

# ============================ 9 - FOUR TIERS ================================
s = new_slide(prs)
y = head(s, 4, "Modeling", "Four tiers, one experiment")
tiers = [("TIER 1", "Majority class.", "always “licit” — the baseline to beat and report"),
         ("TIER 2", "Behaviour only.", "logistic regression & random forest, 29 features"),
         ("TIER 3", "Behaviour + graph.", "the same two models; only the feature set changes"),
         ("TIER 4", "Graph neural network.", "GraphSAGE on the raw address graph — no hand-made features")]
hrule(s, PAD_L, y, CONTENT_W, INK, 1.2)
ry = y
for tn, lead, rest in tiers:
    text(s, PAD_L, ry + 22, 90, tn, size=12, font=MONO, color=CRIMSON, spc=0.8, h=20)
    text(s, PAD_L + 110, ry + 18, CONTENT_W - 120,
         [(lead + " ", {"bold": True}), (rest, {"italic": True, "color": MUTED, "size": 18.5})],
         size=21, h=32)
    ry += 62
    hrule(s, PAD_L, ry, CONTENT_W, HAIR2, 1)
text(s, PAD_L, ry + 30, CONTENT_W,
     "CLASS-WEIGHTED · 5-FOLD CV ON THE TRAINING SPLIT ONLY · ONE TEMPORAL SPLIT THROUGHOUT",
     size=12, font=MONO, color=MUTED, spc=0.5, h=24)

# ============================ 10 - RESULTS ==================================
s = new_slide(prs)
y = head(s, 5, "Results", "Nothing changes but the feature set")
for i, (lab, col) in enumerate([("behaviour", BLUE_LT), ("+ graph features", BLUE), ("raw graph", TEAL)]):
    lx = PAD_L + i * 200
    rect(s, lx, y + 4, 12, 12, fill=col)
    text(s, lx + 20, y, 180, lab, size=11.5, font=MONO, color=MUTED, spc=0.5, h=18)
bot = barchart(s, PAD_L, y + 40, 700, [
    ("logistic regression", [("behaviour", .737, BLUE_LT), ("+ graph", .812, BLUE)]),
    ("random forest",       [("behaviour", .881, BLUE_LT), ("+ graph", .890, BLUE)]),
    ("GraphSAGE",           [("raw graph", .831, TEAL)])])
text(s, PAD_L, bot + 44, 360, "▲ tier 1 baseline · F₁ = 0.000", size=11, font=MONO, color=CRIMSON, h=18)
kstats(s, PAD_L + 760, y + 40, 328,
       [("F₁ = 0.890", "shipped · RF + graph", BLUE),
        ("129 of 133", "scams caught on test", INK),
        ("28", "false alarms · the review queue", CRIMSON)])
text(s, PAD_L, bot + 78, 700,
     [("Fig. 4", {"bold": True, "color": INK}),
      (" — F₁ on the illicit class, held-out temporal test split (856 addresses, 133 illicit).", {})],
     size=11.5, font=MONO, color=MUTED, line=1.75, h=40)

# ============================ 11 - DOES THE GRAPH HELP ======================
s = new_slide(prs)
y = head(s, 5, "The question", "Does the graph help?",
         dek="Yes — but modestly, and it depends on the model.", dl=1)
trio = [("+0.075", "F₁ · logistic regression\na linear model cannot build this itself", BLUE),
        ("+0.009", "F₁ · random forest\n+0.017 PR-AUC", INK),
        ("≈9%", "of forest importance\nout_degree, community_size, pagerank", INK)]
for i, (v, k, col) in enumerate(trio):
    x = PAD_L + i * (CONTENT_W / 3)
    if i: seg(s, x - 16, y + 10, x - 16, y + 168, HAIR, 1)
    text(s, x, y + 20, CONTENT_W / 3 - 40, v, size=62, bold=True, color=col, h=76)
    text(s, x, y + 120, CONTENT_W / 3 - 40, k, size=11, font=MONO, color=MUTED,
         spc=1.2, line=1.7, h=60, caps=True)
text(s, PAD_L, y + 212, 880,
     "We would rather report +0.009 accurately than round it into a story. "
     "Our graph is one hop deep — this is a floor, not a ceiling.",
     size=20, italic=True, color=MUTED, line=1.45, h=70)

# ============================ 12 - THE GNN ==================================
s = new_slide(prs)
y = head(s, 5, "Graph neural network", "What if the model reads\nthe graph itself?", tl=2, tw=640)
text(s, PAD_L, y, 560,
     [("A two-layer ", {}), ("GraphSAGE", {"bold": True, "color": INK, "italic": False}),
      (" over all 44,432 nodes — including unlabelled neighbours — given ", {}),
      ("no", {"italic": True}),
      (" graph features at all. It has to rediscover what network context means by "
       "passing messages along the edges.", {})],
     size=20, italic=True, color=MUTED, line=1.45, h=140)
text(s, PAD_L, y + 172, 560,
     [("SAME TEMPORAL SPLIT · SCALER ON TRAINING ROWS ONLY\n", {}),
      ("neighbour_risk_ratio WITHHELD — FEEDING IT OUR ANSWER WOULD MAKE THE COMPARISON CIRCULAR.", {"color": INK})],
     size=12, font=MONO, color=MUTED, spc=0.5, line=1.8, h=70)
kstats(s, PAD_L + 700, y - 10, 388,
       [("F₁ = 0.831", "tier 4 · GraphSAGE", TEAL),
        ("0.884", "PR-AUC", INK),
        ("128 of 133", "scams caught · 96% recall", INK)])
text(s, PAD_L, y + 286, 660,
     [("It loses to the forest on precision — 1,996 labelled training nodes is a small "
       "corpus for learned features. ", {}),
      ("That the forest still wins is evidence our feature engineering wasn’t leaving much behind.",
       {"bold": True, "color": INK, "italic": False})],
     size=20, italic=True, color=MUTED, line=1.45, h=80)

# ============================ 13 - SENSITIVITY ==============================
s = new_slide(prs)
y = head(s, 6, "Evaluation", "Is the signal real?",
         dek="Licit addresses were sampled from random blocks, so busy addresses are "
             "over-picked: 61% hit the 100-transaction cap, against 8% of illicit ones.",
         dl=2, dw=840)
bot = barchart(s, PAD_L + 60, y + 46, 900, [
    ("n_tx alone",                [("a", .836, BLUE_LT)]),
    ("activity features removed", [("b", .887, BLUE_LT)]),
    ("activity-matched pairs",    [("c", .892, BLUE)])], marker=(.890, "headline 0.890"), LBL=210)
text(s, PAD_L, bot + 60, CONTENT_W,
     [("Fig. 5", {"bold": True, "color": INK}),
      (" — within activity-matched pairs, activity carries no information by construction "
       "— and F₁ does not drop. ", {}),
      ("The confound is genuine but not load-bearing.", {"bold": True, "color": INK})],
     size=11.5, font=MONO, color=MUTED, align=PP_ALIGN.CENTER, line=1.75, h=44)

# ============================ 14 - CODE STRUCTURE ===========================
s = new_slide(prs)
y = head(s, 7, "Code", "One path from raw\nchain to application", tl=2, tw=640)
nodes = [("labels.py", "darklist + era matching", "solid"),
         ("collect.py", "resumable, rate-limit aware", "solid"),
         ("transactions\n.jsonl.gz", "frozen — every stage reads this", "dash"),
         ("features.py · graph.py", "29 + 7 · leakage guard", "solid"),
         ("train.py · gnn.py", "4 tiers, temporal split", "solid"),
         ("model.joblib", "scaler + classifier together", "dash"),
         ("app/app.py", "streamlit + docker", "ship")]
NW, GAPN = 136, 22
for i, (nm, kk, style) in enumerate(nodes):
    nx = PAD_L + i * (NW + GAPN)
    col = BLUE if style == "ship" else (MUTED if style == "dash" else INK)
    box = rect(s, nx, y, NW, 100, fill=PAPER, line_col=col, line_w=2 if style == "ship" else 1)
    if style == "dash":
        box.line.dash_style = DASH.DASH
    text(s, nx + 6, y + 12, NW - 12, nm, size=13, font=MONO, bold=True,
         align=PP_ALIGN.CENTER, line=1.25, h=36)
    text(s, nx + 6, y + 56, NW - 12, kk, size=9.5, font=MONO, color=MUTED,
         align=PP_ALIGN.CENTER, spc=0.5, line=1.55, h=36, caps=True)
    if i < len(nodes) - 1:
        text(s, nx + NW, y + 40, GAPN, "→", size=13, font=MONO, color=MUTED,
             align=PP_ALIGN.CENTER, h=20)
for i, (hh, body) in enumerate([
        ("NO DUPLICATED PREPROCESSING",
         "The app imports src/features.py rather than reimplementing it — the classic way a "
         "deployed model silently diverges from the one that was evaluated."),
        ("TESTS WHERE BUGS HIDE",
         "Cursor construction, feature arithmetic on hand-built transactions, the leakage "
         "guard, the saved scaler.")]):
    cx = PAD_L + i * (CONTENT_W / 2 + 28)
    text(s, cx, y + 140, CONTENT_W / 2 - 40, hh, size=12, font=MONO, color=CRIMSON, spc=1.4, h=20)
    text(s, cx, y + 168, CONTENT_W / 2 - 40, body, size=17, italic=True, color=MUTED, line=1.45, h=110)

# ============================ 15 - APPLICATION ==============================
s = new_slide(prs)
y = head(s, 8, "Application", "A probability and a priority\n— never a verdict", tl=2, tw=660)
tabs = [("TAB 1", "Look up an address.", "fetched live, scored by the saved pipeline"),
        ("TAB 2", "Explore features by hand.", "presets like “collect-and-sweep”, no network needed"),
        ("TAB 3", "Model performance.", "all four tiers, importances, graph statistics")]
hrule(s, PAD_L, y, 600, INK, 1.2)
ry = y
for tn, lead, rest in tabs:
    text(s, PAD_L, ry + 20, 70, tn, size=12, font=MONO, color=BLUE, spc=0.8, h=20)
    text(s, PAD_L + 84, ry + 16, 516,
         [(lead + " ", {"bold": True}), (rest, {"italic": True, "color": MUTED, "size": 18})],
         size=20, line=1.35, h=62)
    ry += 86
    hrule(s, PAD_L, ry, 600, HAIR2, 1)
BX = PAD_L + 660
rect(s, BX, y + 6, 428, 96, fill=PAPER, line_col=INK, line_w=1)
text(s, BX, y + 34, 428, "eth-illicit-detection.streamlit.app", size=16, font=MONO,
     bold=True, align=PP_ALIGN.CENTER, h=26)
text(s, BX, y + 66, 428, "live now — switching to the app", size=10.5, font=MONO,
     color=MUTED, align=PP_ALIGN.CENTER, spc=1.4, h=18, caps=True)
text(s, BX, y + 132, 428,
     [("The footer states it plainly: a LICIT score means ", {}),
      ("“not reported”", {"italic": True}),
      (", not “cleared”. A flag is a candidate for review, not a suspect.", {})],
     size=17.5, italic=True, color=MUTED, line=1.45, h=120)

# ============================ 16 - TASK DIVISION ============================
s = new_slide(prs)
y = head(s, 9, "Closing", "Task division & further work")
cols = [("Sanskriti", ["Blockscout client & collector", "labels & era-matched sampling",
                       "EDA notebook", "Streamlit app & deployment"]),
        ("Sweta", ["behavioural & graph features", "leakage guard, unit tests",
                   "four-tier training & the GNN", "sensitivity analysis"]),
        ("Further work", ["crawl one hop deeper", "more positives — add OFAC",
                          "recalibrate to the true base rate",
                          "measure rate limits before designing around them"])]
CW = CONTENT_W / 3 - 24
for i, (name, items) in enumerate(cols):
    cx = PAD_L + i * (CONTENT_W / 3)
    text(s, cx, y, CW, name, size=21, bold=True, h=30)
    hrule(s, cx, y + 40, CW, INK, 1)
    iy = y + 58
    for it in items:
        text(s, cx, iy, 12, "·", size=17.5, color=CRIMSON, h=24)
        text(s, cx + 16, iy, CW - 16, it, size=17.5, color=MUTED, line=1.5, h=52)
        iy += 30 if len(it) < 34 else 56

# ============================ 17 - THE ANSWER ===============================
s = new_slide(prs)
text(s, PAD_L, 206, 760, "Can an address’s network reveal fraud its own features cannot?",
     size=11.5, font=MONO, color=MUTED, spc=1.8, h=20, caps=True)
text(s, PAD_L, 250, 880,
     [("Yes", {"italic": True, "color": BLUE}),
      (" — decisive for a linear model, marginal for a forest, and confirmed by a network "
       "that reads the raw graph without beating the features we built by hand.", {})],
     size=42, bold=True, line=1.26, h=290)
text(s, PAD_L, 610, 900, "Sanskriti Acharya & Sweta Sharma · eth-illicit-detection.streamlit.app",
     size=12, font=MONO, color=MUTED, spc=0.5, h=22)

out = sys.argv[1] if len(sys.argv) > 1 else "deck.pptx"
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
